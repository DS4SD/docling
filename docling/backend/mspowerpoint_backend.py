import logging
from io import BytesIO
from pathlib import Path
from typing import Set, Union

from docling_core.types.doc import (
    BoundingBox,
    CoordOrigin,
    DocItemLabel,
    DoclingDocument,
    DocumentOrigin,
    GroupLabel,
    ImageRef,
    ProvenanceItem,
    Size,
    TableCell,
    TableData,
)
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from docling.backend.abstract_backend import (
    DeclarativeDocumentBackend,
    PaginatedDocumentBackend,
)
from docling.datamodel.base_models import InputFormat
from docling.datamodel.document import InputDocument

_log = logging.getLogger(__name__)


class MsPowerpointDocumentBackend(DeclarativeDocumentBackend, PaginatedDocumentBackend):
    def __init__(self, in_doc: "InputDocument", path_or_stream: Union[BytesIO, Path]):
        super().__init__(in_doc, path_or_stream)
        self.namespaces = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        }
        # Powerpoint file:
        self.path_or_stream = path_or_stream

        self.pptx_obj = None
        self.valid = False
        try:
            if isinstance(self.path_or_stream, BytesIO):
                self.pptx_obj = Presentation(self.path_or_stream)
            elif isinstance(self.path_or_stream, Path):
                self.pptx_obj = Presentation(str(self.path_or_stream))

            self.valid = True
        except Exception as e:
            raise RuntimeError(
                f"MsPowerpointDocumentBackend could not load document with hash {self.document_hash}"
            ) from e

        return

    def page_count(self) -> int:
        if self.is_valid():
            assert self.pptx_obj is not None
            return len(self.pptx_obj.slides)
        else:
            return 0

    def is_valid(self) -> bool:
        return self.valid

    @classmethod
    def supports_pagination(cls) -> bool:
        return True  # True? if so, how to handle pages...

    def unload(self):
        if isinstance(self.path_or_stream, BytesIO):
            self.path_or_stream.close()

        self.path_or_stream = None

    @classmethod
    def supported_formats(cls) -> Set[InputFormat]:
        return {InputFormat.PPTX}

    def convert(self) -> DoclingDocument:
        # Parses the PPTX into a structured document model.
        # origin = DocumentOrigin(filename=self.path_or_stream.name, mimetype=next(iter(FormatToMimeType.get(InputFormat.PPTX))), binary_hash=self.document_hash)

        origin = DocumentOrigin(
            filename=self.file.name or "file",
            mimetype="application/vnd.ms-powerpoint",
            binary_hash=self.document_hash,
        )

        doc = DoclingDocument(
            name=self.file.stem or "file", origin=origin
        )  # must add origin information
        doc = self.walk_linear(self.pptx_obj, doc)

        return doc

    def generate_prov(
        self, shape, slide_ind, text="", slide_size=Size(width=1, height=1)
    ):
        if shape.left:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
        else:
            left = 0
            top = 0
            width = slide_size.width
            height = slide_size.height
        shape_bbox = [left, top, left + width, top + height]
        shape_bbox = BoundingBox.from_tuple(shape_bbox, origin=CoordOrigin.BOTTOMLEFT)
        prov = ProvenanceItem(
            page_no=slide_ind + 1, charspan=[0, len(text)], bbox=shape_bbox
        )

        return prov

    def handle_text_elements(self, shape, parent_slide, slide_ind, doc, slide_size):
        is_a_list = False
        is_list_group_created = False
        enum_list_item_value = 0
        new_list = None
        bullet_type = "None"
        list_text = ""
        list_label = GroupLabel.LIST
        doc_label = DocItemLabel.LIST_ITEM
        prov = self.generate_prov(shape, slide_ind, shape.text.strip(), slide_size)

        # Identify if shape contains lists
        for paragraph in shape.text_frame.paragraphs:
            # Check if paragraph is a bullet point using the `element` XML
            p = paragraph._element
            if (
                p.find(".//a:buChar", namespaces={"a": self.namespaces["a"]})
                is not None
            ):
                bullet_type = "Bullet"
                is_a_list = True
            elif (
                p.find(".//a:buAutoNum", namespaces={"a": self.namespaces["a"]})
                is not None
            ):
                bullet_type = "Numbered"
                is_a_list = True
            else:
                is_a_list = False

            if paragraph.level > 0:
                # Most likely a sub-list
                is_a_list = True

            if is_a_list:
                # Determine if this is an unordered list or an ordered list.
                # Set GroupLabel.ORDERED_LIST when it fits.
                if bullet_type == "Numbered":
                    list_label = GroupLabel.ORDERED_LIST

            if is_a_list:
                _log.debug("LIST DETECTED!")
            else:
                _log.debug("No List")

        # If there is a list inside of the shape, create a new docling list to assign list items to
        # if is_a_list:
        #     new_list = doc.add_group(
        #         label=list_label, name=f"list", parent=parent_slide
        #     )

        # Iterate through paragraphs to build up text
        for paragraph in shape.text_frame.paragraphs:
            # p_text = paragraph.text.strip()
            p = paragraph._element
            enum_list_item_value += 1
            inline_paragraph_text = ""
            inline_list_item_text = ""

            for e in p.iterfind(".//a:r", namespaces={"a": self.namespaces["a"]}):
                if len(e.text.strip()) > 0:
                    e_is_a_list_item = False
                    is_numbered = False
                    if (
                        p.find(".//a:buChar", namespaces={"a": self.namespaces["a"]})
                        is not None
                    ):
                        bullet_type = "Bullet"
                        e_is_a_list_item = True
                    elif (
                        p.find(".//a:buAutoNum", namespaces={"a": self.namespaces["a"]})
                        is not None
                    ):
                        bullet_type = "Numbered"
                        is_numbered = True
                        e_is_a_list_item = True
                    else:
                        e_is_a_list_item = False

                    if e_is_a_list_item:
                        if len(inline_paragraph_text) > 0:
                            # output accumulated inline text:
                            doc.add_text(
                                label=doc_label,
                                parent=parent_slide,
                                text=inline_paragraph_text,
                                prov=prov,
                            )
                        # Set marker and enumerated arguments if this is an enumeration element.
                        inline_list_item_text += e.text
                        # print(e.text)
                    else:
                        # Assign proper label to the text, depending if it's a Title or Section Header
                        # For other types of text, assign - PARAGRAPH
                        doc_label = DocItemLabel.PARAGRAPH
                        if shape.is_placeholder:
                            placeholder_type = shape.placeholder_format.type
                            if placeholder_type in [
                                PP_PLACEHOLDER.CENTER_TITLE,
                                PP_PLACEHOLDER.TITLE,
                            ]:
                                # It's a title
                                doc_label = DocItemLabel.TITLE
                            elif placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                                DocItemLabel.SECTION_HEADER
                        enum_list_item_value = 0
                        inline_paragraph_text += e.text

            if len(inline_paragraph_text) > 0:
                # output accumulated inline text:
                doc.add_text(
                    label=doc_label,
                    parent=parent_slide,
                    text=inline_paragraph_text,
                    prov=prov,
                )

            if len(inline_list_item_text) > 0:
                enum_marker = ""
                if is_numbered:
                    enum_marker = str(enum_list_item_value) + "."
                if not is_list_group_created:
                    new_list = doc.add_group(
                        label=list_label, name=f"list", parent=parent_slide
                    )
                    is_list_group_created = True
                doc.add_list_item(
                    marker=enum_marker,
                    enumerated=is_numbered,
                    parent=new_list,
                    text=inline_list_item_text,
                    prov=prov,
                )
        return

    def handle_title(self, shape, parent_slide, slide_ind, doc):
        placeholder_type = shape.placeholder_format.type
        txt = shape.text.strip()
        prov = self.generate_prov(shape, slide_ind, txt)

        if len(txt.strip()) > 0:
            # title = slide.shapes.title.text if slide.shapes.title else "No title"
            if placeholder_type in [PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.TITLE]:
                _log.info(f"Title found: {shape.text}")
                doc.add_text(
                    label=DocItemLabel.TITLE, parent=parent_slide, text=txt, prov=prov
                )
            elif placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                _log.info(f"Subtitle found: {shape.text}")
                # Using DocItemLabel.FOOTNOTE, while SUBTITLE label is not avail.
                doc.add_text(
                    label=DocItemLabel.SECTION_HEADER,
                    parent=parent_slide,
                    text=txt,
                    prov=prov,
                )
        return

    def handle_pictures(self, shape, parent_slide, slide_ind, doc, slide_size):
        # Open it with PIL
        try:
            # Get the image bytes
            image = shape.image
            image_bytes = image.blob
            im_dpi, _ = image.dpi
            pil_image = Image.open(BytesIO(image_bytes))

            # shape has picture
            prov = self.generate_prov(shape, slide_ind, "", slide_size)
            doc.add_picture(
                parent=parent_slide,
                image=ImageRef.from_pil(image=pil_image, dpi=im_dpi),
                caption=None,
                prov=prov,
            )
        except (UnidentifiedImageError, OSError) as e:
            _log.warning(f"Warning: image cannot be loaded by Pillow: {e}")
        return

    def handle_tables(self, shape, parent_slide, slide_ind, doc, slide_size):
        # Handling tables, images, charts
        if shape.has_table:
            table = shape.table
            table_xml = shape._element

            prov = self.generate_prov(shape, slide_ind, "", slide_size)

            num_cols = 0
            num_rows = len(table.rows)
            tcells = []
            # Access the XML element for the shape that contains the table
            table_xml = shape._element

            for row_idx, row in enumerate(table.rows):
                if len(row.cells) > num_cols:
                    num_cols = len(row.cells)
                for col_idx, cell in enumerate(row.cells):
                    # Access the XML of the cell (this is the 'tc' element in table XML)
                    cell_xml = table_xml.xpath(
                        f".//a:tbl/a:tr[{row_idx + 1}]/a:tc[{col_idx + 1}]"
                    )

                    if not cell_xml:
                        continue  # If no cell XML is found, skip

                    cell_xml = cell_xml[0]  # Get the first matching XML node
                    row_span = cell_xml.get("rowSpan")  # Vertical span
                    col_span = cell_xml.get("gridSpan")  # Horizontal span

                    if row_span is None:
                        row_span = 1
                    else:
                        row_span = int(row_span)

                    if col_span is None:
                        col_span = 1
                    else:
                        col_span = int(col_span)

                    icell = TableCell(
                        text=cell.text.strip(),
                        row_span=row_span,
                        col_span=col_span,
                        start_row_offset_idx=row_idx,
                        end_row_offset_idx=row_idx + row_span,
                        start_col_offset_idx=col_idx,
                        end_col_offset_idx=col_idx + col_span,
                        col_header=False,
                        row_header=False,
                    )
                    if len(cell.text.strip()) > 0:
                        tcells.append(icell)
            # Initialize Docling TableData
            data = TableData(num_rows=num_rows, num_cols=num_cols, table_cells=[])
            # Populate
            for tcell in tcells:
                data.table_cells.append(tcell)
            if len(tcells) > 0:
                # If table is not fully empty...
                # Create Docling table
                doc.add_table(parent=parent_slide, data=data, prov=prov)
        return

    def walk_linear(self, pptx_obj, doc) -> DoclingDocument:
        # Units of size in PPTX by default are EMU units (English Metric Units)
        slide_width = pptx_obj.slide_width
        slide_height = pptx_obj.slide_height

        text_content = []  # type: ignore

        max_levels = 10
        parents = {}  # type: ignore
        for i in range(0, max_levels):
            parents[i] = None

        # Loop through each slide
        for slide_num, slide in enumerate(pptx_obj.slides):
            slide_ind = pptx_obj.slides.index(slide)
            parent_slide = doc.add_group(
                name=f"slide-{slide_ind}", label=GroupLabel.CHAPTER, parent=parents[0]
            )

            slide_size = Size(width=slide_width, height=slide_height)
            parent_page = doc.add_page(page_no=slide_ind + 1, size=slide_size)

            def handle_shapes(shape, parent_slide, slide_ind, doc, slide_size):
                handle_groups(shape, parent_slide, slide_ind, doc, slide_size)
                if shape.has_table:
                    # Handle Tables
                    self.handle_tables(shape, parent_slide, slide_ind, doc, slide_size)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Handle Pictures
                    self.handle_pictures(
                        shape, parent_slide, slide_ind, doc, slide_size
                    )
                # If shape doesn't have any text, move on to the next shape
                if not hasattr(shape, "text"):
                    return
                if shape.text is None:
                    return
                if len(shape.text.strip()) == 0:
                    return
                if not shape.has_text_frame:
                    _log.warning("Warning: shape has text but not text_frame")
                    return
                # Handle other text elements, including lists (bullet lists, numbered lists)
                self.handle_text_elements(
                    shape, parent_slide, slide_ind, doc, slide_size
                )
                return

            def handle_groups(shape, parent_slide, slide_ind, doc, slide_size):
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for groupedshape in shape.shapes:
                        handle_shapes(
                            groupedshape, parent_slide, slide_ind, doc, slide_size
                        )

            # Loop through each shape in the slide
            for shape in slide.shapes:
                handle_shapes(shape, parent_slide, slide_ind, doc, slide_size)

        return doc
